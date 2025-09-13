# backend/extract_text.py
import sys, os, boto3, io, json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
from dotenv import load_dotenv

load_dotenv()

# Set Tesseract CMD path - check PATH first, then try default location
if os.name == 'nt':  # Windows
    import shutil
    
    # Try to find tesseract in PATH first
    tesseract_path = shutil.which('tesseract')
    if tesseract_path:
        print(f"Found Tesseract in PATH: {tesseract_path}", file=sys.stderr)
        pytesseract.pytesseract.tesseract_cmd = tesseract_path
    else:
        # Fall back to default location
        default_path = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        if os.path.exists(default_path):
            print(f"Using default Tesseract path: {default_path}", file=sys.stderr)
            pytesseract.pytesseract.tesseract_cmd = default_path
        else:
            print("WARNING: Tesseract not found in PATH or default location. OCR will not work.", file=sys.stderr)

def check_tesseract():
    """Verify Tesseract installation and configuration."""
    try:
        # Try a simple OCR operation to verify Tesseract works
        test_image = Image.new('RGB', (100, 30), color='white')
        pytesseract.image_to_string(test_image)
        return True
    except Exception as e:
        print(json.dumps({"error": f"Tesseract not properly configured: {str(e)}. Please ensure Tesseract is installed."}), file=sys.stderr)
        return False

def perform_ocr(image_bytes):
    """Performs OCR on image bytes and returns extracted text."""
    try:
        image = Image.open(io.BytesIO(image_bytes))
        
        # Convert RGBA to RGB if needed (Tesseract prefers RGB)
        if image.mode == 'RGBA':
            background = Image.new('RGB', image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[3])
            image = background
        
        # Attempt to improve image quality for OCR
        image = image.convert('L')  # Convert to grayscale
        
        # Run OCR with specific configuration for better results
        custom_config = r'--oem 3 --psm 6'  # Page segmentation mode 6: Assume uniform block of text
        text = pytesseract.image_to_string(image, config=custom_config)
        extracted = text.strip()
        
        if extracted:
            print(f"OCR succeeded, extracted {len(extracted)} characters", file=sys.stderr)
            return extracted
        else:
            print("OCR produced no text", file=sys.stderr)
            return ""
    except Exception as e:
        print(f"OCR Error: {e}", file=sys.stderr)
        return ""

def download_from_s3(bucket_name, s3_key):
    try:
        s3 = boto3.client('s3')
        buf = io.BytesIO()
        s3.download_fileobj(bucket_name, s3_key, buf)
        buf.seek(0)
        return buf
    except Exception as e:
        print(json.dumps({"error": f"S3 download failed: {e}"}), file=sys.stderr)
        sys.exit(1)

def count_slides(file_path_or_buffer):
    """Just count the slides in a PowerPoint file without extracting text."""
    try:
        # Improved debugging
        print(f"Attempting to count slides in: {file_path_or_buffer}", file=sys.stderr)
        
        # If it's a string path, verify file exists
        if isinstance(file_path_or_buffer, str):
            if not os.path.exists(file_path_or_buffer):
                print(f"File does not exist: {file_path_or_buffer}", file=sys.stderr)
                return 0
            print(f"File exists at {file_path_or_buffer}, size: {os.path.getsize(file_path_or_buffer)} bytes", file=sys.stderr)
        
        prs = Presentation(file_path_or_buffer)
        slide_count = len(prs.slides)
        print(f"Successfully counted {slide_count} slides", file=sys.stderr)
        return slide_count
    except Exception as e:
        print(json.dumps({"error": f"Error counting slides: {str(e)}"}), file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        return 0  # Return 0 instead of exiting to handle errors better

def extract_text_from_presentation(file_buffer, count_only=False):
    """Extract text and/or count slides from a PowerPoint presentation.
    
    Args:
        file_buffer: File object or path to PowerPoint file
        count_only: If True, only count slides without extracting text
        
    Returns:
        If count_only is True: Dictionary with slideCount
        Otherwise: List of slide data with extracted text
    """
    if count_only:
        slide_count = count_slides(file_buffer)
        return {"slideCount": slide_count}
        
    data = []
    try:
        # First verify Tesseract is working if we need OCR (skip for count_only)
        if not count_only and not check_tesseract():
            return data

        prs = Presentation(file_buffer)
        if not prs.slides: return data

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            slide_texts = []
            notes_text = None
            image_count = 0
            ocr_count = 0

            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    try:
                        image_bytes = shape.image.blob
                        ocr_text = perform_ocr(image_bytes)
                        if ocr_text:
                            ocr_count += 1
                            slide_texts.append(f"[OCR Text: {ocr_text}]")
                            print(f"Slide {slide_num}: OCR successful for image {image_count}", file=sys.stderr)
                    except Exception as img_err:
                        print(f"Error processing image {image_count} on slide {slide_num}: {img_err}", file=sys.stderr)

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            # Include OCR statistics in the slide data
            slide_data = {
                "slide": slide_num,
                "text": " ".join(slide_texts),
                "notes": notes_text,
                "stats": {
                    "total_images": image_count,
                    "ocr_successful": ocr_count
                }
            }
            data.append(slide_data)
            
            # Log slide processing summary
            print(f"Processed slide {slide_num}: {image_count} images, {ocr_count} OCR successes", file=sys.stderr)

        return data
    except Exception as e:
        print(json.dumps({"error": f"PPTX processing failed: {e}"}), file=sys.stderr)
        return data  # Return empty data instead of exiting to handle errors better

# Export function for Node.js integration
def extractText(file_path, count_only=False):
    """Function to be called from Node.js."""
    try:
        print(f"extractText called with file_path={file_path}, count_only={count_only}", file=sys.stderr)
        result = extract_text_from_presentation(file_path, count_only)
        return result
    except Exception as e:
        print(f"Error in extractText: {e}", file=sys.stderr)
        return {"error": str(e)}

if __name__ == "__main__":
    # Handle both direct file paths and S3 keys
    if len(sys.argv) < 2:
        print(json.dumps({"error": "Usage: python extract_text.py <file_path_or_s3_key> [count_only]"}))
        sys.exit(1)

    input_path = sys.argv[1]
    count_only = len(sys.argv) > 2 and sys.argv[2].lower() == 'count_only'
    
    # Determine if it's a file path or S3 key
    if os.path.isfile(input_path):
        # It's a local file
        print(f"Processing local file: {input_path}", file=sys.stderr)
        result = extract_text_from_presentation(input_path, count_only)
    else:
        # Assume it's an S3 key
        print(f"Processing S3 key: {input_path}", file=sys.stderr)
        bucket_name = os.getenv("S3_BUCKET_NAME")
        if not bucket_name:
            print(json.dumps({"error": "S3_BUCKET_NAME environment variable not set"}))
            sys.exit(1)
        
        presentation_buffer = download_from_s3(bucket_name, input_path)
        result = extract_text_from_presentation(presentation_buffer, count_only)
    
    # Output the result
    print(json.dumps({"data": result}))